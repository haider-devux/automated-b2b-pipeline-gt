"""Generate briefing.pptx — a full project deck (16:9, dark theme). Run with the venv Python."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0E, 0x11, 0x16)
PURPLE = RGBColor(0x7C, 0x5C, 0xFF)
WHITE = RGBColor(0xF3, 0xF5, 0xF8)
MUTED = RGBColor(0x9A, 0xA6, 0xB5)
GREEN = RGBColor(0x36, 0xD0, 0x7A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def txt(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = "Calibri"
    return tb


def title_slide():
    s = prs.slides.add_slide(BLANK); bg(s, PURPLE)
    txt(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.6),
        "Granjur B2B Acquisition Pipeline", 44, WHITE, bold=True)
    txt(s, Inches(1.05), Inches(3.75), Inches(11), Inches(0.8),
        "Automated lead generation & outreach - built in-house, on free tools", 20, WHITE)
    rect(s, Inches(1.05), Inches(4.7), Inches(2.2), Inches(0.06), WHITE)
    txt(s, Inches(1.05), Inches(6.4), Inches(11), Inches(0.5), "Project briefing", 14, WHITE)


def content(title, kicker, lines, foot=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    rect(s, Inches(0.6), Inches(0.72), Inches(0.12), Inches(0.85), PURPLE)
    txt(s, Inches(0.9), Inches(0.6), Inches(11), Inches(0.4), kicker, 14, PURPLE, bold=True)
    txt(s, Inches(0.85), Inches(1.0), Inches(11.8), Inches(1.0), title, 32, WHITE, bold=True)
    tb = s.shapes.add_textbox(Inches(0.95), Inches(2.35), Inches(11.4), Inches(4.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (head, sub) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = "▪  " + head
        r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        if sub:
            r2 = p.add_run(); r2.text = "    " + sub
            r2.font.size = Pt(15); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"
    if foot:
        txt(s, Inches(0.95), Inches(6.75), Inches(11.4), Inches(0.5), foot, 13, MUTED, italic=True)


title_slide()

content("The goal", "WHY WE BUILT IT", [
    ("Granjur needs a steady flow of qualified clients.", ""),
    ("Doing it by hand is slow.", "finding, researching and writing to each lead takes hours."),
    ("And it's risky.", "a few bad emails can get our domains blocked everywhere."),
    ("So we built a system that does the heavy lifting - and keeps us safe.", ""),
])

content("How it works", "THE BIG PICTURE", [
    ("1.  Discover", "find companies from free public sources"),
    ("2.  Enrich", "research each one automatically"),
    ("3.  Qualify + Pitch", "AI decides fit and writes a tailored message"),
    ("4.  Outreach", "auto-sends - or a person can approve first"),
], foot="ONE command runs all four stages. A central database is the brain; a live dashboard is the window.")

content("WF-1   -   Discovery", "FIND THE COMPANIES", [
    ("Free public sources", "online maps data (OpenStreetMap) + developer job boards"),
    ("Humans add leads too", "team pastes in companies spotted on LinkedIn / Upwork"),
    ("Auto-filtered", "keeps only the right size and region; skips big chains"),
], foot="Tools: Python, OpenStreetMap, RemoteOK & Remotive job feeds")

content("WF-2   -   Enrichment", "RESEARCH EACH ONE", [
    ("Tech detection", "spots what a company's website is built with"),
    ("Contact finding", "pulls a real public email from their site"),
    ("Verification", "checks the email address actually works"),
], foot="Tools: Python, website analysis, email / DNS verification - all free")

content("WF-3   -   AI Qualification & Pitch", "DECIDE FIT + WRITE THE MESSAGE", [
    ("Fit decision by rules, not guesswork", "right size + right signals -> keep or drop"),
    ("Segment", "matches each lead to one of three Granjur offers"),
    ("AI-written pitch", "personalised, in English, Arabic or Chinese"),
], foot="Tools: local AI (Ollama) on our own PC + a deterministic rules engine")

content("WF-4   -   Outreach", "SEND & TRACK", [
    ("Fully automatic, or human-checked", "runs hands-off - or you approve pitches on the dashboard first"),
    ("Routed to the email tool", "each message mapped to the right campaign and sending domain"),
    ("Tracks the result", "replies, bookings and unsubscribes flow back in"),
], foot="Recommended: keep the human check until real sending is trusted. In practice mode - nothing emails yet.")

content("The dashboard", "THE WINDOW INTO EVERYTHING", [
    ("Funnel & leads", "see every company and where it is in the process"),
    ("Review queue", "approve or reject discovered companies and pitches"),
    ("Analytics", "conversion and yield by source and segment"),
], foot="A simple web page on our own PC - no logins, no subscriptions")

content("The technology", "A 100% FREE STACK", [
    ("Python", "powers all four components"),
    ("PostgreSQL database", "the single source of truth"),
    ("Local AI (Ollama)", "writes the pitches, offline, no API bills"),
    ("Free data sources", "OpenStreetMap, RemoteOK, Remotive, public websites"),
], foot="No paid subscriptions or API fees anywhere in the pipeline")

content("Compliance & safety", "PROTECTING OUR DOMAINS", [
    ("Unsubscribe + suppression list", "opt-outs are honoured permanently"),
    ("No role-account emails", "avoids spam-trap addresses like info@ / sales@"),
    ("Region rules", "stricter countries go to manual, not cold email"),
    ("Human approval + practice mode", "nothing sends without a person - and nothing sends yet"),
])

content("Status & what's next", "WHERE WE ARE", [
    ("Done", "all four stages run from ONE command, free, in practice mode"),
    ("Next", "warm up sending domains, connect the email tool, go live"),
    ("Then", "feed results back in to sharpen targeting over time"),
], foot="From an n8n prototype to a complete, free, in-house pipeline.")

out = os.path.join(os.path.dirname(__file__), "briefing.pptx")
prs.save(out)
print("wrote", out, "-", len(prs.slides._sldIdLst), "slides")
